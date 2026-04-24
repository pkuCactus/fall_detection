from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
DARK_BG = RGBColor(30, 30, 40)
ACCENT = RGBColor(70, 130, 180)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
ORANGE = RGBColor(230, 126, 34)
GREEN = RGBColor(46, 204, 113)

def add_slide_with_bg(prs, title_text, subtitle_text=None):
    """添加带背景色的空白幻灯片"""
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)
    # 添加背景色矩形
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    # 副标题
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY
    return slide

def add_bullet_text(slide, text, left=Inches(0.5), top=Inches(1.8), width=Inches(12.3), height=Inches(5.2), font_size=Pt(16)):
    """添加带项目符号的文本框"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
        if line.strip().startswith('•') or line.strip().startswith('-') or line.strip().startswith('|'):
            p.level = 1
        elif line.strip().startswith('  ') and (line.strip()[2:].startswith('•') or line.strip()[2:].startswith('-')):
            p.level = 2
    return box

# ============ Slide 1: 封面 ============
slide = add_slide_with_bg(prs, "Edge-AI 跌倒检测系统", "双 Pipeline 架构设计详解")
box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.5))
tf = box.text_frame
tf.word_wrap = True
lines = [
    "目标平台: HiSilicon 3516C (0.5T INT8, 15M DDR, 30M Flash)",
    "核心约束: 端到端 latency ~34ms，纯视觉 pipeline",
    "",
    "Pipeline A: 精细化多模型融合路线",
    "Pipeline B: 端到端开放词汇轻量路线"
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_GRAY
    p.space_after = Pt(12)

# ============ Slide 2: 目录 ============
slide = add_slide_with_bg(prs, "目录")
text = """
1. 系统总览：双 Pipeline 架构对比
2. Pipeline A — 主 Pipeline (FallDetectionPipeline)
   • 检测 → 跟踪 → 姿态估计 → 分类器 → 规则引擎 → 融合决策
   • 抽帧优化与关键点缓存
3. Pipeline B — YOLO-World Pipeline (YOLOWorldFallPipeline)
   • 开放词汇检测 → 姿态映射 → 运动辅助 → 融合决策
4. 核心模块深度解析
5. 配置体系与部署策略
6. 总结与选型建议
"""
add_bullet_text(slide, text, top=Inches(1.8))

# ============ Slide 3: 系统总览 ============
slide = add_slide_with_bg(prs, "系统总览 — 双 Pipeline 架构")
text = """
┌─────────────────────────────┬───────────────────────────────────────────┐
│     Pipeline A (主)          │        Pipeline B (YOLO-World)            │
│  FallDetectionPipeline       │      YOLOWorldFallPipeline                │
├─────────────────────────────┼───────────────────────────────────────────┤
│  YOLOv8n 人体检测            │  YOLO-World 开放词汇多姿态检测              │
│       ↓                      │       ↓                                   │
│  ByteTrack-lite 跟踪         │  ByteTrack-lite 跟踪                      │
│       ↓                      │       ↓                                   │
│  YOLOv8n-pose 关键点估计      │  (无独立姿态估计)                          │
│       ↓                      │       ↓                                   │
│  3分支融合分类器 / 简单图像分类 │  姿态类别 → 跌倒得分映射                    │
│       ↓                      │       ↓                                   │
│  RuleEngine 6条规则          │  运动特征 + 宽高比辅助                      │
│       ↓                      │       ↓                                   │
│  FusionDecision 状态机       │  FusionDecision 状态机                    │
│       ↓                      │       ↓                                   │
│      告警输出                 │      告警输出                              │
└─────────────────────────────┴───────────────────────────────────────────┘
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 4: Pipeline A 总流程 ============
slide = add_slide_with_bg(prs, "Pipeline A 总流程", "FallDetectionPipeline — 端到端数据流")
text = """
每帧输入 (BGR, HWC)
    │
    ├── 检测帧? ──Yes──► YOLOv8n PersonDetector ──► Detection[]
    │                       │
    │                       ▼
    │               ByteTrackLite.update()
    │                       │
    │                       ▼
    │               YOLOv8n-pose (整图推理 + IoU匹配)
    │                       │
    │                       ▼
    │               FallClassifier / SimpleFallClassifier
    │                       │
    │                       ▼
    │               RuleEngine.evaluate(kpts, bbox, history, cls_score)
    │                       │
    │                       ▼
    │               FusionDecision.update(rule_score, cls_score, posture)
    │
    └── 跳过帧? ──No───► Tracker.predict() only
                            │
                            ▼
                    复用上帧 kpts / 关键点跟踪器预测
                            │
                            ▼
                    Classifier 每帧仍推理 (用缓存/预测 kpts)
                            │
                            ▼
                    RuleEngine + FusionDecision (用新 cls_score)

关键设计: 重模型（检测、姿态）只在检测帧运行；分类器、规则、融合每帧都运行，保证时序连续性。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 5: 检测层 ============
slide = add_slide_with_bg(prs, "Pipeline A — 检测层", "PersonDetector (src/fall_detection/core/detector.py)")
text = """
底层模型: YOLOv8n (默认) / 自定义训练权重
输入: np.ndarray HWC, BGR (OpenCV 默认)
输出: List[Dict{bbox, conf, class_id, class_name}]
推理尺寸: imgsz 可配置 (默认 640)
后处理: conf_thresh 过滤 + filter_class_id 类别筛选

代码要点:
  class PersonDetector:
      def __init__(self, model_name="yolov8n", model_path=None,
                   classes=None, device=None, model_type="yolo", imgsz=None)

支持 model_type="yolo" 和 "yolo_world" 两种后端
调用 self.model(img, verbose=False, imgsz=self.imgsz)
遍历 result.boxes，提取 xyxy, conf, cls
"""
add_bullet_text(slide, text, top=Inches(1.8))

# ============ Slide 6: 跟踪层 ============
slide = add_slide_with_bg(prs, "Pipeline A — 跟踪层", "ByteTrackLite (src/fall_detection/core/tracker.py)")
text = """
设计目标: 裁剪版 ByteTrack，仅用 IoU + Kalman，无 ReID，极致轻量。

Step 1: 所有已有 Track.predict() ──► KalmanFilter 预测下一帧位置
Step 2: 分离高/低分检测框
    ├── dets_high: conf >= track_thresh (默认 0.5)
    └── dets_low:  conf < track_thresh
Step 3: 第一次匹配 ── confirmed/tentative tracks ↔ dets_high
    代价矩阵: 1 - IoU，匈牙利算法，阈值 match_thresh (默认 0.8)
Step 4: 第二次匹配 ── 未匹配 tracks ↔ dets_low (低分框可能是遮挡物体)
Step 5: 未匹配 track → mark_missed()
    tentative track 直接删除
    confirmed track 允许 time_since_update <= max_age (默认 30)
Step 6: 未匹配高分检测框 → 新建 tentative Track
    hits >= min_hits (默认 3) 后升级为 confirmed

KalmanFilter: 8 状态 (x, y, w, h, vx, vy, vw, vh)，匀速模型，Cholesky 分解求解增益。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 7: 姿态估计 ============
slide = add_slide_with_bg(prs, "Pipeline A — 姿态估计", "PoseEstimator (src/fall_detection/core/pose_estimator.py)")
text = """
设计权衡: 整图推理一次，再按 IoU 关联到检测框 —— 避免逐 crop 推理的重复开销。

输入: frame (BGR), bboxes[]
    │
    ▼
YOLOv8n-pose 整图推理
    │
    ▼
提取 pose_boxes[] (xyxy) + pose_kpts[] (17, 3)
    过滤 class_id != 0 (只保留 person)
    │
    ▼
对每个输入 bbox，贪心匹配最佳 pose_box:
    ├── 遍历所有未使用的 pose_box
    ├── 计算 IoU，取最大
    └── 若 best_iou > 0.1 → 关联成功，返回对应 kpts
        否则 → 返回 zeros((17, 3))

输出: List[np.ndarray]，每个 (17, 3) = [x, y, conf]

注意点: 遮挡或密集人群时，IoU 关联可能出错；未匹配框返回全零关键点。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 8: 分类器 ============
slide = add_slide_with_bg(prs, "Pipeline A — 分类器设计", "FallClassifier vs SimpleFallClassifier")
text = """
方案一: FallClassifier — 3分支融合分类器

输入 1: ROI 图像 (3, 96, 96)
    img_conv1(3→16, s=2) → img_conv2(16→32, s=2) → GAP → 32-d
输入 2: 17 关键点 (17, 3) → flatten 51-d
    kpt_fc(51→32) → ReLU → 32-d
输入 3: 8-d 运动特征 [vx, vy, ax, ay, w, h, h_ratio, n_ground]
    motion_fc(8→8) → ReLU → 8-d
Concat(32+32+8=72) → fusion_fc1(72→32) → Dropout(0.3) → fusion_fc2(32→1) → Sigmoid

方案二: SimpleFallClassifier — 单分支图像分类器 (边缘优先)

Backbone:
  Conv(3→64, s=2) → BN → ReLU
  Conv(64→64, s=2) → BN → ReLU
  ResBlock(64→64) → ReLU → BasicBlock(64→64) → ReLU
  ResBlock(64→128, s=2) → ReLU → BasicBlock(128→128) → AvgPool → Flatten → 128-d
Head:
  Linear(128→128) → BN → Dropout → ReLU → Linear(128→2) → Softmax

Pipeline 自适应: cfg.classifier.type 为 "simple" 时，自动降低融合权重 alpha、提高 beta、降低 alarm_thresh。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 9: 规则引擎 ============
slide = add_slide_with_bg(prs, "Pipeline A — 规则引擎", "RuleEngine 6条规则 + 姿态分类 (src/fall_detection/core/rules.py)")
text = """
输入: kpts(17,3), bbox, history{centers[]}, cls_score
输出: score[0,1], flags{A,B,C,D,E,F}, debug{posture, h_ratio, ...}

6 条规则详解:
A 高度压缩比: h_ratio 低 + 地面接触点 ≥ 2 —— 核心跌倒特征：身体压缩贴地
B 地面区域: 若配置 ground_roi，检查 bbox 底部是否进入地面区域；否则 fallback 检查底边贴近 bbox 底部
C 动→静转换: 早期位移 avg_early > motion_thresh 且 晚期 avg_late < static_thresh；或 lying + 全程静止
D 垂直快速下降: vy_px_s > fall_vy_thresh + (h_ratio 低 或 lying)
E 向下加速度: accel > accel_thresh + (h_ratio 低 或 lying)
F 关键点不可见: 可见关键点比例 < no_keypoint_thresh —— 极端姿态/遮挡时兜底

姿态分类逻辑 (_classify_posture):
  kpt_aspect > 1.5 or torso_angle > 55          → lying  (光轴方向跌倒)
  bbox_aspect > 1.3 and kpt_aspect > 1.0 and h_ratio < 0.65 → lying
  h_ratio > 0.75                                 → standing
  h_ratio > 0.50:
      hip_ratio < 0.45  → sitting
      else              → standing (if h_ratio>0.65) else sitting
  h_ratio > 0.35:
      hip_ratio < 0.45  → sitting   (减少坐姿误判)
      else              → crouching
  else                                             → lying

分类器辅助: cls_score > t1(0.85) 强制 lying；cls_score < t2(0.2) 强制 standing；中间区间 standing 降级为 sitting。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(12))

# ============ Slide 10: 融合决策 ============
slide = add_slide_with_bg(prs, "Pipeline A — 融合决策", "FusionDecision 五态状态机 (src/fall_detection/core/fusion.py)")
text = """
不是简单加权求和！而是驱动 FallState 状态机。

状态机流转:
NORMAL ──(连续N帧超阈值)──► SUSPECTED ──(序列检查/分类器bypass)──► FALLING
                                                              │
                                                              ▼
                                                        ALARM_SENT ──(should_alarm=True)
                                                              │
                                        (score低于阈值持续alarm_reset_frames)
                                                              ▼
                                                        RECOVERING ──(持续recovery_frames)
                                                              │
                                                              ▼
                                                           NORMAL

分数融合公式:
  S_final   = alpha * s_rule + beta * s_cls + gamma * S_temporal
  S_temporal = EMA / 滑动窗口平滑 (历史分类器分数)
  默认权重: alpha=0.5, beta=0.3, gamma=0.2

关键约束:
  alarm_thresh (0.5)           — 进入 suspected 的阈值
  alarm_min_frames (3)         — suspected 状态持续帧数才允许进入 falling
  sequence_check_frames (8)    — 姿态历史检查窗口：必须有 upright(站/坐) + fall(蹲/躺)，且 fall_count > upright_count
  cls_bypass_thresh (0.85)     — 分类器极高置信度时，可绕过姿态序列检查，且 alarm_min_frames 减半
  cooldown_seconds (3s)        — 告警冷却期，防止重复上报
  recovery_seconds (0.5s)      — 恢复确认期，快速恢复正常
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 11: 抽帧优化 ============
slide = add_slide_with_bg(prs, "Pipeline A — 抽帧优化与缓存策略", "核心设计: skip_frames = 2 → 每 3 帧做一次检测")
text = """
目标: 满足 edge latency 预算 (~34ms 总耗时, ≤15M DDR)

检测帧 (0, 3, 6, ...):
  ├─ 跑 YOLOv8n 检测      (~10ms)
  ├─ ByteTrack update      (~1ms)
  ├─ YOLOv8n-pose 整图推理 (~8ms)
  ├─ Classifier 推理       (~3ms)
  ├─ RuleEngine            (~1ms)
  └─ FusionDecision        (~0.5ms)

跳过帧 (1, 2, 4, 5, ...):
  ├─ Tracker.predict()     (~0.2ms)
  ├─ 复用/预测关键点
  ├─ Classifier 推理       (~3ms)   ← 每帧都跑，保证 cls_score 连续
  ├─ RuleEngine            (~1ms)
  └─ FusionDecision        (~0.5ms)

缓存机制:
  _last_track_kpts    — 跳过帧复用上一帧关键点，track 死亡时移除
  _last_cls_scores    — 检测帧保存，跳过帧复用，每检测帧更新
  _track_history      — 中心点历史 (deque, maxlen=1.5s)，自动淘汰旧帧
  _fusion_score_history — 融合分数时序 (可视化用，deque maxlen=30)

关键点跟踪器 (SimpleKeypointTracker, 默认禁用):
  enabled=false 时：跳帧直接复用上帧 kpts（简单但可能闪烁）
  enabled=true  时：用 smooth_alpha + velocity_decay 预测插值，可选光流辅助
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 12: Pipeline B 总流程 ============
slide = add_slide_with_bg(prs, "Pipeline B 总流程", "YOLOWorldFallPipeline — 开放词汇端到端")
text = """
每帧输入 (BGR, HWC)
    │
    ├── 检测帧? ──Yes──► YOLO-World (开放词汇)
    │   │               直接输出 8 类姿态: standing, sitting, squatting,
    │   │               bending, half_up_or_crouching, kneeling,
    │   │               crawling_or_crawling_like, lying_on_floor
    │   │                   │
    │   │                   ▼
    │   │           ByteTrackLite.update()
    │   │                   │
    │   │                   ▼
    │   │           IoU 匹配: 检测类别 → Track
    │   │           (保存 det_info: class_name, class_id, conf)
    │   │
    └── 跳过帧? ──No───► Tracker.predict() only
                            复用 _last_det_info 的类别信息
    │
    ▼
_update_track_history()   (记录中心点轨迹)
    │
    ▼
_process_tracks()
    │
    ├── _compute_fall_score()
    │   ├── base_score = fall_scores[class_name]   (配置映射)
    │   ├── motion_bonus = 垂直快速下降加分
    │   └── aspect_bonus = 宽高比 > 1.2 加分
    │   └── score = base_score * det_conf + motion_bonus*0.1 + aspect_bonus
    │
    ├── FusionDecision.update(rule_score=score, cls_score=score, posture=posture)
    │   (复用同一套状态机！)
    │
    └── 输出: tracks, scores, falling_flags, alarms
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 13: YOLO-World 检测 ============
slide = add_slide_with_bg(prs, "Pipeline B — YOLO-World 开放词汇检测", "检测即姿态")
text = """
传统方案: 检测器只输出 "person" → 后接 pose estimator → 后接规则/分类器推断姿态。
YOLO-World 方案: 检测器直接输出细粒度姿态类别。

classes:
  - "person standing"
  - "person sitting"
  - "person squatting"
  - "person bending"
  - "person half up or crouching"
  - "person kneeling"
  - "person crawling or crawling-like"
  - "person lying on floor"

优势:
  • 无需独立的 pose estimator 模型 (省 ~8ms, 省内存)
  • 无需关键点规则引擎 (省 ~1ms)
  • 端到端训练，姿态判别由检测器内部完成

劣势:
  • 依赖开放词汇模型的零样本/少样本能力
  • 姿态类别固定，难以微调 (相比关键点规则的灵活性)
  • 类别混淆风险 (squatting vs crouching vs kneeling)
"""
add_bullet_text(slide, text, top=Inches(1.8))

# ============ Slide 14: 姿态映射与跌倒得分 ============
slide = add_slide_with_bg(prs, "Pipeline B — 姿态映射与跌倒得分", "两层映射体系")
text = """
Layer 1: 检测类别 → 标准姿态 (posture_map)

  person standing              → standing
  person sitting               → sitting
  person squatting             → crouching
  person bending               → crouching
  person half up or crouching  → crouching
  person kneeling              → crouching
  person crawling or crawling-like → lying
  person lying on floor        → lying

目的: 将 8 个开放词汇类别归一化为 4 个标准姿态，供 FusionDecision 状态机使用。

Layer 2: 检测类别 → 基础跌倒分 (fall_scores)

  person standing              : 0.00
  person sitting               : 0.15
  person squatting             : 0.25
  person bending               : 0.35
  person half up or crouching  : 0.45
  person kneeling              : 0.30
  person crawling or crawling-like : 0.55
  person lying on floor        : 0.95

最终得分计算:
  score = base_score * det_conf + motion_bonus * 0.1 + aspect_bonus

  det_conf: 检测置信度作为基础分权重 (高置信度 → 高分更可信)
  motion_bonus: 垂直下降速度 > 50px 时加分 (max 0.5)
  aspect_bonus: bbox 宽高比 > 1.2 时加分 (max 0.2)
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 15: 运动与外观辅助 ============
slide = add_slide_with_bg(prs, "Pipeline B — 运动与外观辅助特征", "弥补无关键点信息的不足")
text = """
1. 运动辅助 (_compute_motion_bonus)

  centers = track_history[tid]  # 最近 0.5s 中心点
  vy = centers[-1].y - centers[-2].y   # 垂直位移 (y 向下为正)
  if vy > 50:                          # 快速下降
      motion_bonus = min(0.5, vy / 400.0)

  作用: 捕捉 "站立 → 突然倒下" 的动态过程，即使检测器类别尚未切换到 lying。

2. 宽高比辅助 (aspect_bonus)

  aspect = w / max(h, 1.0)
  if aspect > 1.2:                     # 倒下的人通常 wider than tall
      aspect_bonus = min(0.2, (aspect - 1.2) * 0.3)

  作用: lying 姿态的几何特征补偿。

3. FusionDecision 状态机 (完全复用)

  # YOLO-World Pipeline 中: rule_score 与 cls_score 均使用 fall_score
  self.fusion[tid].update(
      rule_score=fall_score,
      cls_score=fall_score,
      posture=posture
  )

  设计意图: 虽然信息来源不同，但时序融合、状态机防抖、姿态序列检查的逻辑完全一致。
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 16: 双 Pipeline 深度对比 ============
slide = add_slide_with_bg(prs, "双 Pipeline 深度对比")
text = """
| 维度               | Pipeline A (主)              | Pipeline B (YOLO-World)      |
|--------------------|------------------------------|------------------------------|
| 核心思想            | 多模型精细化融合              | 端到端开放词汇轻量            |
| 检测器              | YOLOv8n (闭集, person)       | YOLO-World (开放词汇, 8姿态)  |
| 姿态估计            | YOLOv8n-pose (17关键点)       | 无 (检测器直接输出姿态)        |
| 分类器              | 3分支融合 / Simple图像分类     | 无 (得分映射替代)             |
| 规则引擎            | 6条规则 + 姿态分类 + cls辅助   | 无 (运动/外观辅助替代)         |
| 融合决策            | FusionDecision 状态机         | 完全复用同一状态机            |
| 跟踪器              | ByteTrackLite                | 同一 ByteTrackLite           |
| 模型数量            | 3~4 个 (det, pose, cls)      | 1 个 (yolo_world)            |
| 抽帧策略            | skip_frames=2                | skip_frames=2                |
| 内存占用            | 较高 (多模型权重)             | 较低 (单模型)                |
| latency             | ~34ms (检测帧)                | ~20ms (检测帧, 估计)          |
| 可解释性            | 高 (关键点可视化, 规则flag)    | 中 (类别+得分可视化)          |
| 场景适配            | 已训练数据的封闭场景           | 零样本/少样本开放场景          |
| 误报来源            | 规则阈值敏感, 遮挡关联错        | 类别混淆, 词汇歧义            |
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 17: 配置体系 ============
slide = add_slide_with_bg(prs, "配置体系", "三层配置分离")
text = """
configs/
├── pipeline/
│   ├── default.yaml          # Pipeline A 配置
│   └── yolo_world.yaml       # Pipeline B 配置
└── training/
    ├── detector.yaml
    ├── pose.yaml
    ├── classifier.yaml
    ├── simple_classifier.yaml
    └── yoloworld.yaml

default.yaml 关键配置段:
  detector:      { conf_thresh: 0.3, model_path: ..., imgsz: null }
  tracker:       { track_thresh: 0.5, match_thresh: 0.8, max_age: 30, min_hits: 3 }
  pose_estimator:{ model_path: ... }
  classifier:    { type: "simple"/"fusion", model_path: ..., fall_class_idx: 1 }
  rules:         { h_ratio_thresh: 0.6, motion_thresh: 50, fall_vy_thresh: 200,
                   cls_posture_t1: 0.85, cls_posture_t2: 0.2, ... }
  fusion:        { alpha: 0.5, beta: 0.3, gamma: 0.2, alarm_thresh: 0.5,
                   alarm_min_frames: 3, sequence_check_frames: 8,
                   cls_bypass_thresh: 0.85, cooldown_seconds: 3.0 }
  keypoint_tracker: { enabled: false, smooth_alpha: 0.7, use_optical_flow: false }
  pipeline:      { skip_frames: 2, fps: 25 }

Pipeline 自适应逻辑:
  if use_simple_classifier:
      fusion_cfg["alpha"] *= 0.3      # 规则权重降低
      fusion_cfg["beta"]  *= 2.0      # 分类器权重提高
      fusion_cfg["alarm_thresh"] *= 0.8  # 降低告警阈值
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 18: 部署与性能 ============
slide = add_slide_with_bg(prs, "部署与性能")
text = """
运行时命令:

  # Pipeline A 演示
  PYTHONPATH=src python scripts/demo/run_pipeline_demo.py \
      --video data/video/test.mp4 \
      --config configs/pipeline/default.yaml

  # Pipeline B 演示
  PYTHONPATH=src python scripts/demo/run_pipeline_demo.py \
      --video data/video/test.mp4 \
      --config configs/pipeline/yolo_world.yaml

  # 速度基准测试
  PYTHONPATH=src python scripts/eval/benchmark_speed.py \
      --video data/video/test.mp4

模型加载路径约定:
  models/pretrained/        # 预训练权重
  outputs/detector/         # 训练产出: 检测器
  outputs/pose/             # 训练产出: 姿态估计器
  outputs/classifier/       # 训练产出: 3分支融合分类器
  outputs/simple_classifier/# 训练产出: 单分支图像分类器

Edge 部署要点:
  • skip_frames=2 是 latency 预算的核心杠杆
  • SimpleFallClassifier 比 FallClassifier 更适合 INT8 量化 (单分支, 无关键点输入)
  • keypoint_tracker.enabled=false 是默认保守策略，避免跳帧时关键点漂移
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(14))

# ============ Slide 19: 总结与选型建议 ============
slide = add_slide_with_bg(prs, "总结与选型建议")
text = """
设计亮点:

1. 统一跟踪与融合层: 两个 Pipeline 复用同一 ByteTrackLite 和 FusionDecision，保证升级时一致性。
2. 抽帧架构: 检测/姿态仅在 1/3 帧运行，分类器每帧运行，平衡精度与性能。
3. 双分类器适配: 3分支融合 (精度优先) vs Simple图像 (速度/量化优先)，Pipeline 层自动调参。
4. 状态机防抖: NORMAL→SUSPECTED→FALLING→ALARM_SENT→RECOVERING 五态机 + cooldown/recovery，杜绝抖动。
5. 开放词汇探索: YOLO-World Pipeline 证明了 "检测即姿态" 的可行性，为边缘极端资源场景提供备选。

选型建议:

| 场景                        | 推荐 Pipeline              | 理由                          |
|-----------------------------|----------------------------|-------------------------------|
| 已标注数据充分，精度优先      | A + FallClassifier         | 关键点规则可解释性强，融合精度高 |
| 边缘 INT8 部署，latency 极敏感 | A + SimpleFallClassifier   | 单分支易量化，模型体积最小       |
| 零样本/开放类别/快速 PoC      | B (YOLO-World)             | 无需训练 pose/cls，单模型即插即用 |
| 混合场景 (已知+未知姿态)       | A 为主，B 为辅              | A 处理常规，B 兜底开放词汇       |
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(13))

# ============ Slide 20: Q&A / 源码索引 ============
slide = add_slide_with_bg(prs, "Q&A — 关键源码索引")
text = """
| 模块             | 文件路径                                              |
|------------------|-------------------------------------------------------|
| 主 Pipeline       | src/fall_detection/pipeline/pipeline.py               |
| YOLO-World Pipeline | src/fall_detection/pipeline/yoloworld_pipeline.py    |
| 检测器            | src/fall_detection/core/detector.py                   |
| 跟踪器            | src/fall_detection/core/tracker.py                    |
| 姿态估计          | src/fall_detection/core/pose_estimator.py             |
| 规则引擎          | src/fall_detection/core/rules.py                      |
| 融合决策          | src/fall_detection/core/fusion.py                     |
| 融合分类器        | src/fall_detection/models/classifier.py               |
| 简单分类器        | src/fall_detection/models/simple_classifier.py        |
| 主配置            | configs/pipeline/default.yaml                         |
| YW 配置           | configs/pipeline/yolo_world.yaml                      |

感谢聆听！
"""
add_bullet_text(slide, text, top=Inches(1.8), font_size=Pt(16))

# 保存
output_path = "/data4/hjz/fall_detection/docs/design/pipeline_architecture.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
